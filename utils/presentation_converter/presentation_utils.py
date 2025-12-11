# utils/presentation_utils.py - PRESENTATION CONVERTER UTILITIES
# Comprehensive presentation conversion utilities following project requirements
import os
import subprocess
import tempfile
import shutil
from pathlib import Path
import logging
from typing import List, Optional, Tuple, Dict, Any

# Universal Security Framework Integration
try:
    from security.core.validators import validate_content, validate_filename
    from security.core.sanitizers import sanitize_filename
    SECURITY_FRAMEWORK_AVAILABLE = True
except ImportError:
    SECURITY_FRAMEWORK_AVAILABLE = False

try:
    from pdf2image import convert_from_path
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.dml.color import RGBColor
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    import pytesseract
    from PIL import Image, ImageEnhance, ImageFilter
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

try:
    import cv2
    import numpy as np
    OPENCV_AVAILABLE = True
except ImportError:
    OPENCV_AVAILABLE = False

# Initialize logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class PresentationConverter:
    """
    Comprehensive presentation converter supporting:
    - PPTX to PDF conversion (using LibreOffice)
    - PDF to PPTX basic conversion (image-based slides)
    - PDF to PPTX accurate conversion (OCR + layout reconstruction)
    """
    
    def __init__(self):
        self.temp_dirs = []  # Track temporary directories for cleanup
        
        # Check for required dependencies
        self.libreoffice_path = self._find_libreoffice()
        self.dependencies = self._check_dependencies()
        
    def __del__(self):
        """Cleanup temporary directories on object destruction"""
        self.cleanup_temp_dirs()
        
    def cleanup_temp_dirs(self):
        """Clean up all temporary directories created by this instance"""
        for temp_dir in self.temp_dirs:
            if os.path.exists(temp_dir):
                try:
                    shutil.rmtree(temp_dir)
                except Exception as e:
                    logger.warning(f"Failed to cleanup temp dir {temp_dir}: {e}")
        self.temp_dirs.clear()
    
    def _find_libreoffice(self) -> Optional[str]:
        """Find LibreOffice installation path"""
        import platform
        
        possible_paths = []
        
        if platform.system() == "Windows":
            # Windows paths - check both 32-bit and 64-bit locations
            possible_paths.extend([
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                "soffice",  # Try PATH
                "libreoffice"  # Try PATH
            ])
        else:
            # Linux/macOS paths
            possible_paths.extend([
                "/usr/bin/libreoffice",
                "/usr/bin/soffice",
                "/snap/bin/libreoffice",
                "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                "soffice",
                "libreoffice"
            ])
        
        for path in possible_paths:
            logger.debug(f"Checking LibreOffice path: {path}")
            try:
                if path in ["soffice", "libreoffice"]:
                    # Enhanced Windows subprocess to prevent console window
                    if platform.system() == "Windows":
                        startupinfo = subprocess.STARTUPINFO()
                        startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW
                        startupinfo.wShowWindow = subprocess.SW_HIDE
                        creation_flags = subprocess.CREATE_NO_WINDOW | subprocess.DETACHED_PROCESS
                    else:
                        startupinfo = None
                        creation_flags = 0
                    
                    result = subprocess.run(
                        [path, "--version"], 
                        capture_output=True, 
                        text=True, 
                        timeout=10,
                        creationflags=creation_flags,
                        startupinfo=startupinfo
                    )
                    # LibreOffice returns 0 for --version if successful
                    if result.returncode == 0:
                        logger.info(f"Found LibreOffice in PATH: {path}")
                        return path
                else:
                    # Check if file exists
                    logger.debug(f"Checking if file exists: {path}")
                    if os.path.exists(path) and os.path.isfile(path):
                        logger.debug(f"File exists: {path}")
                        # For Windows, if the file exists and is in the LibreOffice directory,
                        # assume it's valid to avoid timeout issues
                        if platform.system() == "Windows" and "LibreOffice" in path:
                            logger.info(f"Found LibreOffice at: {path}")
                            return path
                        else:
                            # For non-Windows or non-standard paths, test execution
                            try:
                                # Enhanced Windows subprocess to prevent console window
                                if platform.system() == "Windows":
                                    startupinfo = subprocess.STARTUPINFO()
                                    startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW
                                    startupinfo.wShowWindow = subprocess.SW_HIDE
                                    creation_flags = subprocess.CREATE_NO_WINDOW | subprocess.DETACHED_PROCESS
                                else:
                                    startupinfo = None
                                    creation_flags = 0
                                
                                result = subprocess.run(
                                    [path, "--version"], 
                                    capture_output=True, 
                                    text=True, 
                                    timeout=5,
                                    creationflags=creation_flags,
                                    startupinfo=startupinfo
                                )
                                if result.returncode == 0:
                                    logger.info(f"Found LibreOffice at: {path}")
                                    return path
                            except (subprocess.TimeoutExpired, subprocess.SubprocessError, OSError) as e:
                                logger.debug(f"Failed to test LibreOffice at {path}: {e}")
                                continue
            except (subprocess.TimeoutExpired, subprocess.SubprocessError, OSError) as e:
                logger.debug(f"Failed to check LibreOffice path {path}: {e}")
                continue
                
        logger.warning("LibreOffice not found in any expected location")
        return None
    
    def _check_dependencies(self) -> Dict[str, bool]:
        """Check availability of all dependencies"""
        return {
            'libreoffice': self.libreoffice_path is not None,
            'pdf2image': PDF2IMAGE_AVAILABLE,
            'python_pptx': PYTHON_PPTX_AVAILABLE,
            'pymupdf': PYMUPDF_AVAILABLE,
            'ocr': OCR_AVAILABLE,
            'opencv': OPENCV_AVAILABLE
        }
    
    def is_pptx_conversion_available(self) -> bool:
        """Check if PPTX to PDF conversion is available"""
        return self.dependencies['libreoffice']
    
    def is_pdf_conversion_available(self) -> bool:
        """Check if basic PDF to PPTX conversion is available"""
        return (self.dependencies['pdf2image'] and 
                self.dependencies['python_pptx'])
    
    def is_ocr_available(self) -> bool:
        """Check if OCR-based accurate conversion is available"""
        return (self.dependencies['ocr'] and 
                self.dependencies['pymupdf'] and 
                self.dependencies['python_pptx'])
    
    def validate_file_security(self, file_path: str) -> Tuple[bool, List[str]]:
        """Validate file using universal security framework"""
        if not SECURITY_FRAMEWORK_AVAILABLE:
            logger.warning("Security framework not available - performing basic validation")
            return self._basic_file_validation(file_path)
        
        try:
            # Read file content for security validation
            with open(file_path, 'rb') as f:
                file_content = f.read()
            
            # Determine file type
            file_ext = Path(file_path).suffix.lower().lstrip('.')
            
            # Perform security validation
            is_safe, issues = validate_content(file_content, file_ext)
            
            if not is_safe:
                logger.error(f"Security validation failed for {file_path}: {issues}")
                return False, issues
            
            logger.info(f"Security validation passed for {file_path}")
            return True, []
            
        except Exception as e:
            logger.error(f"Security validation error for {file_path}: {e}")
            return False, [f"Security validation error: {str(e)}"]
    
    def _basic_file_validation(self, file_path: str) -> Tuple[bool, List[str]]:
        """Basic file validation when security framework is not available"""
        try:
            if not os.path.exists(file_path):
                return False, ["File does not exist"]
            
            file_size = os.path.getsize(file_path)
            if file_size == 0:
                return False, ["File is empty"]
            
            if file_size > 100 * 1024 * 1024:  # 100MB limit
                return False, ["File too large"]
            
            return True, []
            
        except Exception as e:
            return False, [f"Basic validation error: {str(e)}"]
    
    def pptx_to_pdf(self, input_path: str, output_path: str, 
                   quality: str = 'high', slide_range: Optional[List[int]] = None) -> bool:
        """
        Convert PowerPoint to PDF using LibreOffice
        
        Args:
            input_path: Path to input PPTX file
            output_path: Path for output PDF file
            quality: 'high', 'medium', or 'low'
            slide_range: List of slide numbers to convert (1-indexed)
            
        Returns:
            bool: True if conversion successful
        """
        try:
            if not self.is_pptx_conversion_available():
                logger.error("LibreOffice not available for PPTX to PDF conversion")
                return False
            
            if not os.path.exists(input_path):
                logger.error(f"Input file not found: {input_path}")
                return False
            
            # Security validation using universal security framework
            is_safe, security_issues = self.validate_file_security(input_path)
            if not is_safe:
                logger.error(f"Security validation failed for {input_path}: {security_issues}")
                return False
            
            # Create temporary directory for this conversion
            temp_dir = tempfile.mkdtemp(prefix='pptx_pdf_')
            self.temp_dirs.append(temp_dir)
            
            input_for_conversion = input_path
            
            # Handle slide range by creating a temporary PPTX with selected slides
            if slide_range:
                input_for_conversion = self._extract_slides(input_path, slide_range, temp_dir)
                if not input_for_conversion:
                    logger.error("Failed to extract specified slide range")
                    return False
            
            # Create Windows-compatible user installation directory
            import platform
            if platform.system() == "Windows":
                user_install_dir = os.path.join(temp_dir, "libreoffice_user")
                os.makedirs(user_install_dir, exist_ok=True)
                user_install_url = f"file:///{user_install_dir.replace(os.sep, '/')}"
            else:
                user_install_dir = "/tmp/libreoffice_tmp"
                user_install_url = "file:///tmp/libreoffice_tmp"
            
            # Prepare LibreOffice command
            cmd = [
                self.libreoffice_path,
                "--headless",
                "--invisible",
                "--nodefault",
                "--nolockcheck",
                "--nologo",
                "--norestore",
                "--convert-to", "pdf",
                "--outdir", os.path.dirname(output_path),
                f"-env:UserInstallation={user_install_url}",
                input_for_conversion
            ]
            
            # Add quality-specific parameters
            # Note: Removed --writer as it's for Writer documents, not Impress
            
            # Execute conversion
            logger.info(f"Converting PPTX to PDF: {input_path} -> {output_path}")
            logger.debug(f"LibreOffice command: {' '.join(cmd)}")
            
            try:
                # Enhanced Windows subprocess flags to prevent console window
                if platform.system() == "Windows":
                    import subprocess
                    startupinfo = subprocess.STARTUPINFO()
                    startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW
                    startupinfo.wShowWindow = subprocess.SW_HIDE
                    creation_flags = subprocess.CREATE_NO_WINDOW | subprocess.DETACHED_PROCESS
                else:
                    startupinfo = None
                    creation_flags = 0
                
                result = subprocess.run(
                    cmd,
                    capture_output=True,
                    text=True,
                    timeout=300,  # 5 minute timeout
                    cwd=temp_dir,
                    creationflags=creation_flags,
                    startupinfo=startupinfo
                )
                
                logger.debug(f"LibreOffice return code: {result.returncode}")
                logger.debug(f"LibreOffice stdout: {result.stdout}")
                logger.debug(f"LibreOffice stderr: {result.stderr}")
                
                if result.returncode != 0:
                    error_msg = result.stderr or result.stdout or "Unknown error"
                    logger.error(f"LibreOffice conversion failed with return code {result.returncode}: {error_msg}")
                    return False
                    
            except subprocess.TimeoutExpired:
                logger.error("LibreOffice conversion timed out after 5 minutes")
                return False
            except Exception as e:
                logger.error(f"Failed to execute LibreOffice: {e}")
                return False
            
            # Find the generated PDF file
            base_name = os.path.splitext(os.path.basename(input_for_conversion))[0]
            generated_pdf = os.path.join(os.path.dirname(output_path), f"{base_name}.pdf")
            
            # Move to final output path if different
            if generated_pdf != output_path:
                if os.path.exists(generated_pdf):
                    shutil.move(generated_pdf, output_path)
                else:
                    logger.error(f"Generated PDF not found at expected location: {generated_pdf}")
                    return False
            
            success = os.path.exists(output_path)
            if success:
                logger.info(f"PPTX to PDF conversion successful: {output_path}")
            else:
                logger.error("PPTX to PDF conversion failed - output file not created")
                
            return success
            
        except subprocess.TimeoutExpired:
            logger.error("LibreOffice conversion timed out")
            return False
        except Exception as e:
            logger.error(f"PPTX to PDF conversion error: {e}")
            return False
    
    def _extract_slides(self, input_path: str, slide_numbers: List[int], temp_dir: str) -> Optional[str]:
        """Extract specific slides from PPTX file"""
        try:
            if not PYTHON_PPTX_AVAILABLE:
                logger.error("python-pptx not available for slide extraction")
                return None
            
            prs = Presentation(input_path)
            new_prs = Presentation()
            
            # Clear default slide
            new_prs.slides._sldIdLst.clear()
            
            # Add selected slides (convert to 0-indexed)
            for slide_num in slide_numbers:
                if 1 <= slide_num <= len(prs.slides):
                    slide_idx = slide_num - 1
                    slide = prs.slides[slide_idx]
                    
                    # Copy slide layout and content
                    slide_layout = new_prs.slide_layouts[0]  # Use default layout
                    new_slide = new_prs.slides.add_slide(slide_layout)
                    
                    # Copy all shapes from original slide
                    for shape in slide.shapes:
                        # Note: This is a simplified copy - full shape copying
                        # would require more complex implementation
                        if hasattr(shape, 'text'):
                            try:
                                new_shape = new_slide.shapes.add_textbox(
                                    shape.left, shape.top, shape.width, shape.height
                                )
                                new_shape.text = shape.text
                            except:
                                pass  # Skip shapes that can't be copied
            
            # Save temporary PPTX
            temp_pptx_path = os.path.join(temp_dir, "extracted_slides.pptx")
            new_prs.save(temp_pptx_path)
            
            return temp_pptx_path
            
        except Exception as e:
            logger.error(f"Slide extraction error: {e}")
            return None
    
    
    

# Convenience function for quick access
def convert_presentation(input_path: str, output_path: str, **kwargs) -> bool:
    """
    Convert PowerPoint to PDF format
    
    Args:
        input_path: Input PPTX/PPT file path
        output_path: Output PDF file path
        **kwargs: Additional conversion parameters
        
    Returns:
        bool: True if successful
    """
    converter = PresentationConverter()
    
    input_ext = Path(input_path).suffix.lower()
    output_ext = Path(output_path).suffix.lower()
    
    if input_ext in ['.pptx', '.ppt'] and output_ext == '.pdf':
        return converter.pptx_to_pdf(input_path, output_path, **kwargs)
    else:
        logger.error(f"Unsupported conversion: {input_ext} -> {output_ext}. Only PPTX/PPT to PDF is supported.")
        return False
