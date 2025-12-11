# utils/pdf_presentation_utils.py - PDF TO PPTX CONVERTER UTILITIES
# Dedicated utilities for PDF to PowerPoint conversion
import os
import tempfile
import shutil
from pathlib import Path
import logging
from typing import List, Optional, Tuple, Dict, Any

# Universal Security Framework Integration
try:
    from security.core.validators import validate_content, validate_filename
    from security.core.sanitizers import sanitize_filename
    SECURITY_FRAMEWORK_AVAILABLE = True
except ImportError:
    SECURITY_FRAMEWORK_AVAILABLE = False

try:
    from pdf2image import convert_from_path
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.dml.color import RGBColor
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    import pytesseract
    from PIL import Image, ImageEnhance, ImageFilter
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

try:
    import cv2
    import numpy as np
    OPENCV_AVAILABLE = True
except ImportError:
    OPENCV_AVAILABLE = False

# Initialize logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class PdfToPptxConverter:
    """
    Dedicated PDF to PowerPoint converter supporting:
    - PDF to PPTX basic conversion (image-based slides)
    - PDF to PPTX accurate conversion (OCR + layout reconstruction)
    """
    
    def __init__(self):
        self.temp_dirs = []  # Track temporary directories for cleanup
        
        # Check for required dependencies
        self.dependencies = self._check_dependencies()
        
    def __del__(self):
        """Cleanup temporary directories on object destruction"""
        self.cleanup_temp_dirs()
        
    def cleanup_temp_dirs(self):
        """Clean up all temporary directories created by this instance"""
        for temp_dir in self.temp_dirs:
            if os.path.exists(temp_dir):
                try:
                    shutil.rmtree(temp_dir)
                except Exception as e:
                    logger.warning(f"Failed to cleanup temp dir {temp_dir}: {e}")
        self.temp_dirs.clear()
    
    def _check_dependencies(self) -> Dict[str, bool]:
        """Check availability of all dependencies"""
        return {
            'pdf2image': PDF2IMAGE_AVAILABLE,
            'python_pptx': PYTHON_PPTX_AVAILABLE,
            'pymupdf': PYMUPDF_AVAILABLE,
            'ocr': OCR_AVAILABLE,
            'opencv': OPENCV_AVAILABLE
        }
    
    def is_pdf_conversion_available(self) -> bool:
        """Check if basic PDF to PPTX conversion is available"""
        return (self.dependencies['pdf2image'] and 
                self.dependencies['python_pptx'])
    
    def is_ocr_available(self) -> bool:
        """Check if OCR-based accurate conversion is available"""
        return (self.dependencies['ocr'] and 
                self.dependencies['pymupdf'] and 
                self.dependencies['python_pptx'])
    
    def validate_file_security(self, file_path: str) -> Tuple[bool, List[str]]:
        """Validate file using universal security framework"""
        if not SECURITY_FRAMEWORK_AVAILABLE:
            logger.warning("Security framework not available - performing basic validation")
            return self._basic_file_validation(file_path)
        
        try:
            # Read file content for security validation
            with open(file_path, 'rb') as f:
                file_content = f.read()
            
            # Determine file type
            file_ext = Path(file_path).suffix.lower().lstrip('.')
            
            # Perform security validation
            is_safe, issues = validate_content(file_content, file_ext)
            
            if not is_safe:
                logger.error(f"Security validation failed for {file_path}: {issues}")
                return False, issues
            
            logger.info(f"Security validation passed for {file_path}")
            return True, []
            
        except Exception as e:
            logger.error(f"Security validation error for {file_path}: {e}")
            return False, [f"Security validation error: {str(e)}"]
    
    def _basic_file_validation(self, file_path: str) -> Tuple[bool, List[str]]:
        """Basic file validation when security framework is not available"""
        try:
            if not os.path.exists(file_path):
                return False, ["File does not exist"]
            
            file_size = os.path.getsize(file_path)
            if file_size == 0:
                return False, ["File is empty"]
            
            if file_size > 100 * 1024 * 1024:  # 100MB limit
                return False, ["File too large"]
            
            return True, []
            
        except Exception as e:
            return False, [f"Basic validation error: {str(e)}"]
    
    def pdf_to_pptx_basic(self, input_path: str, output_path: str) -> bool:
        """
        Convert PDF to PPTX using basic mode (image-based slides)
        
        Args:
            input_path: Path to input PDF file
            output_path: Path for output PPTX file
            
        Returns:
            bool: True if conversion successful
        """
        try:
            if not self.is_pdf_conversion_available():
                logger.error("Required dependencies not available for PDF to PPTX conversion")
                return False
            
            if not os.path.exists(input_path):
                logger.error(f"Input file not found: {input_path}")
                return False
            
            # Security validation using universal security framework
            is_safe, security_issues = self.validate_file_security(input_path)
            if not is_safe:
                logger.error(f"Security validation failed for {input_path}: {security_issues}")
                return False
            
            logger.info(f"Converting PDF to PPTX (Basic): {input_path} -> {output_path}")
            
            # Convert PDF pages to images
            images = convert_from_path(input_path, dpi=200, fmt='PNG')
            
            if not images:
                logger.error("No pages found in PDF")
                return False
            
            # Create new presentation
            prs = Presentation()
            
            # Remove default slide
            prs.slides._sldIdLst.clear()
            
            # Process each page
            for i, img in enumerate(images):
                try:
                    # Add slide with blank layout
                    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
                    slide = prs.slides.add_slide(blank_slide_layout)
                    
                    # Save image temporarily
                    temp_img_path = os.path.join(tempfile.gettempdir(), f"pdf_page_{i}.png")
                    img.save(temp_img_path, 'PNG')
                    
                    # Calculate dimensions to fit slide
                    slide_width = prs.slide_width
                    slide_height = prs.slide_height
                    img_width, img_height = img.size
                    
                    # Calculate scaling to fit while maintaining aspect ratio
                    width_ratio = slide_width / img_width
                    height_ratio = slide_height / img_height
                    scale_ratio = min(width_ratio, height_ratio) * 0.9  # 90% of slide size
                    
                    new_width = int(img_width * scale_ratio)
                    new_height = int(img_height * scale_ratio)
                    
                    # Center the image
                    left = (slide_width - new_width) // 2
                    top = (slide_height - new_height) // 2
                    
                    # Add image to slide
                    slide.shapes.add_picture(temp_img_path, left, top, new_width, new_height)
                    
                    # Clean up temporary image
                    os.unlink(temp_img_path)
                    
                except Exception as e:
                    logger.warning(f"Failed to process page {i+1}: {e}")
                    continue
            
            # Save presentation
            prs.save(output_path)
            
            success = os.path.exists(output_path)
            if success:
                logger.info(f"PDF to PPTX basic conversion successful: {output_path}")
            else:
                logger.error("PDF to PPTX basic conversion failed - output file not created")
                
            return success
            
        except Exception as e:
            logger.error(f"PDF to PPTX basic conversion error: {e}")
            return False
    
    def pdf_to_pptx_accurate(self, input_path: str, output_path: str,
                           ocr_language: str = 'eng', text_detection: str = 'auto',
                           preserve_images: bool = True) -> bool:
        """
        Convert PDF to PPTX using accurate mode (OCR + layout reconstruction)
        
        Args:
            input_path: Path to input PDF file
            output_path: Path for output PPTX file
            ocr_language: OCR language code
            text_detection: 'auto', 'strict', or 'relaxed'
            preserve_images: Whether to preserve images from PDF
            
        Returns:
            bool: True if conversion successful
        """
        try:
            if not self.is_ocr_available():
                logger.error("Required dependencies not available for accurate PDF to PPTX conversion")
                return False
            
            if not os.path.exists(input_path):
                logger.error(f"Input file not found: {input_path}")
                return False
            
            # Security validation using universal security framework
            is_safe, security_issues = self.validate_file_security(input_path)
            if not is_safe:
                logger.error(f"Security validation failed for {input_path}: {security_issues}")
                return False
            
            logger.info(f"Converting PDF to PPTX (Accurate): {input_path} -> {output_path}")
            
            # Open PDF with PyMuPDF
            pdf_doc = fitz.open(input_path)
            
            # Create new presentation
            prs = Presentation()
            prs.slides._sldIdLst.clear()
            
            # Process each page
            for page_num in range(len(pdf_doc)):
                try:
                    page = pdf_doc.load_page(page_num)
                    
                    # Add slide
                    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
                    slide = prs.slides.add_slide(blank_slide_layout)
                    
                    # Extract and process page content
                    self._process_page_accurate(page, slide, ocr_language, 
                                              text_detection, preserve_images)
                    
                except Exception as e:
                    logger.warning(f"Failed to process page {page_num+1} in accurate mode: {e}")
                    continue
            
            pdf_doc.close()
            
            # Save presentation
            prs.save(output_path)
            
            success = os.path.exists(output_path)
            if success:
                logger.info(f"PDF to PPTX accurate conversion successful: {output_path}")
            else:
                logger.error("PDF to PPTX accurate conversion failed - output file not created")
                
            return success
            
        except Exception as e:
            logger.error(f"PDF to PPTX accurate conversion error: {e}")
            return False
    
    def _process_page_accurate(self, page, slide, ocr_language: str, 
                             text_detection: str, preserve_images: bool):
        """Process a single PDF page for accurate conversion"""
        try:
            # Get page dimensions
            page_rect = page.rect
            slide_width = slide.presentation.slide_width
            slide_height = slide.presentation.slide_height
            
            # Calculate scaling factor
            scale_x = slide_width / page_rect.width
            scale_y = slide_height / page_rect.height
            scale = min(scale_x, scale_y) * 0.9
            
            # Extract text blocks with positioning
            text_dict = page.get_text("dict")
            
            for block in text_dict["blocks"]:
                if "lines" in block:  # Text block
                    self._add_text_block(block, slide, scale, page_rect)
                elif preserve_images and "image" in block:  # Image block
                    self._add_image_block(block, slide, scale, page_rect)
            
            # If no text found, fall back to OCR
            if not any("lines" in block for block in text_dict["blocks"]):
                self._perform_ocr_on_page(page, slide, ocr_language, scale, page_rect)
                
        except Exception as e:
            logger.warning(f"Error processing page accurately: {e}")
    
    def _add_text_block(self, block, slide, scale, page_rect):
        """Add a text block to the slide"""
        try:
            # Calculate position and size
            bbox = fitz.Rect(block["bbox"])
            left = int(bbox.x0 * scale)
            top = int(bbox.y0 * scale)
            width = int((bbox.x1 - bbox.x0) * scale)
            height = int((bbox.y1 - bbox.y0) * scale)
            
            # Extract text content
            text_content = ""
            font_size = 12
            
            for line in block["lines"]:
                for span in line["spans"]:
                    text_content += span["text"] + " "
                    if span["size"] > font_size:
                        font_size = span["size"]
                text_content += "\n"
            
            text_content = text_content.strip()
            
            if text_content:
                # Add text box
                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                text_frame.text = text_content
                
                # Set font size
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(max(8, min(int(font_size * scale), 72)))
                        
        except Exception as e:
            logger.warning(f"Error adding text block: {e}")
    
    def _add_image_block(self, block, slide, scale, page_rect):
        """Add an image block to the slide"""
        try:
            # This would require extracting and embedding images from PDF
            # For now, we'll skip image extraction as it's complex
            pass
        except Exception as e:
            logger.warning(f"Error adding image block: {e}")
    
    def _perform_ocr_on_page(self, page, slide, ocr_language, scale, page_rect):
        """Perform OCR on page if no text found"""
        try:
            # Convert page to image
            mat = fitz.Matrix(2, 2)  # 2x zoom for better OCR
            pix = page.get_pixmap(matrix=mat)
            img_data = pix.pil_tobytes(format="PNG")
            
            # Save temporary image
            temp_img_path = os.path.join(tempfile.gettempdir(), "ocr_page.png")
            with open(temp_img_path, "wb") as f:
                f.write(img_data)
            
            # Perform OCR
            img = Image.open(temp_img_path)
            
            # Enhance image for better OCR
            img = img.convert('RGB')
            enhancer = ImageEnhance.Contrast(img)
            img = enhancer.enhance(2.0)
            
            # OCR configuration based on text detection mode
            config = f'--oem 3 --psm 6 -l {ocr_language}'
            if text_detection == 'strict':
                config += ' -c tessedit_char_whitelist=ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz0123456789 '
            
            # Get OCR data with bounding boxes
            ocr_data = pytesseract.image_to_data(img, config=config, output_type=pytesseract.Output.DICT)
            
            # Process OCR results
            self._process_ocr_results(ocr_data, slide, scale)
            
            # Clean up
            os.unlink(temp_img_path)
            
        except Exception as e:
            logger.warning(f"OCR processing error: {e}")
    
    def _process_ocr_results(self, ocr_data, slide, scale):
        """Process OCR results and add text boxes"""
        try:
            current_paragraph = ""
            current_bbox = None
            
            for i in range(len(ocr_data['text'])):
                text = ocr_data['text'][i].strip()
                confidence = int(ocr_data['conf'][i])
                
                if text and confidence > 30:  # Only include high-confidence text
                    x = int(ocr_data['left'][i] * scale / 2)  # Adjust for 2x zoom
                    y = int(ocr_data['top'][i] * scale / 2)
                    w = int(ocr_data['width'][i] * scale / 2)
                    h = int(ocr_data['height'][i] * scale / 2)
                    
                    if current_paragraph:
                        current_paragraph += " " + text
                        # Expand bounding box
                        if current_bbox:
                            current_bbox = [
                                min(current_bbox[0], x),
                                min(current_bbox[1], y),
                                max(current_bbox[2], x + w),
                                max(current_bbox[3], y + h)
                            ]
                    else:
                        current_paragraph = text
                        current_bbox = [x, y, x + w, y + h]
                
                # End of line/paragraph
                if not text or ocr_data['block_num'][i] != ocr_data['block_num'][min(i + 1, len(ocr_data['text']) - 1)]:
                    if current_paragraph and current_bbox:
                        # Add text box
                        textbox = slide.shapes.add_textbox(
                            current_bbox[0], current_bbox[1],
                            current_bbox[2] - current_bbox[0], current_bbox[3] - current_bbox[1]
                        )
                        textbox.text_frame.text = current_paragraph
                        
                        # Set font size
                        for paragraph in textbox.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(max(8, min(14, (current_bbox[3] - current_bbox[1]) // 2)))
                    
                    current_paragraph = ""
                    current_bbox = None
                    
        except Exception as e:
            logger.warning(f"Error processing OCR results: {e}")

# Convenience function for quick access
def convert_pdf_to_pptx(input_path: str, output_path: str, **kwargs) -> bool:
    """
    Convert PDF to PowerPoint format
    
    Args:
        input_path: Input PDF file path
        output_path: Output PPTX file path
        **kwargs: Additional conversion parameters
        
    Returns:
        bool: True if successful
    """
    converter = PdfToPptxConverter()
    
    input_ext = Path(input_path).suffix.lower()
    output_ext = Path(output_path).suffix.lower()
    
    if input_ext == '.pdf' and output_ext in ['.pptx']:
        mode = kwargs.get('mode', 'basic')
        if mode == 'accurate':
            return converter.pdf_to_pptx_accurate(input_path, output_path, **kwargs)
        else:
            return converter.pdf_to_pptx_basic(input_path, output_path)
    else:
        logger.error(f"Unsupported conversion: {input_ext} -> {output_ext}")
        return False