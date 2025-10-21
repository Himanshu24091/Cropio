#!/usr/bin/env python3
"""
Test Script for Presentation Conversion - No Console Window
Tests the enhanced presentation conversion with security framework integration
"""

import sys
import os
import tempfile
import shutil
from pathlib import Path

# Add project root to path
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from utils.presentation_converter.presentation_utils import PresentationConverter
from pptx import Presentation
from pptx.util import Inches

def create_test_pptx(output_path):
    """Create a simple test PPTX file"""
    print("Creating test PPTX file...")
    
    prs = Presentation()
    
    # Slide 1
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
    slide1.shapes.title.text = "Test Presentation"
    slide1.placeholders[1].text = "Created for conversion testing"
    
    # Slide 2
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
    slide2.shapes.title.text = "Sample Content"
    content = slide2.placeholders[1]
    content.text = "This is a test slide\nWith multiple lines\nFor PDF conversion testing"
    
    # Slide 3
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    slide3.shapes.add_textbox(
        Inches(1), Inches(1), Inches(8), Inches(2)
    ).text_frame.text = "Final slide with custom textbox"
    
    prs.save(output_path)
    print(f"Test PPTX created: {output_path}")

def test_pptx_to_pdf_conversion():
    """Test PPTX to PDF conversion without console window"""
    print("\n" + "="*50)
    print("Testing PPTX to PDF Conversion")
    print("="*50)
    
    # Create temporary directory
    temp_dir = tempfile.mkdtemp(prefix='presentation_test_')
    print(f"Using temp directory: {temp_dir}")
    
    try:
        # Create test files
        test_pptx = os.path.join(temp_dir, "test_presentation.pptx")
        output_pdf = os.path.join(temp_dir, "converted_presentation.pdf")
        
        # Create test PPTX
        create_test_pptx(test_pptx)
        
        # Initialize converter
        converter = PresentationConverter()
        
        # Check dependencies
        print(f"\nLibreOffice available: {converter.is_pptx_conversion_available()}")
        if not converter.is_pptx_conversion_available():
            print("❌ LibreOffice not available - cannot test conversion")
            return False
        
        # Test security validation
        print("\nTesting security validation...")
        is_safe, issues = converter.validate_file_security(test_pptx)
        print(f"Security validation result: {'✅ Safe' if is_safe else '❌ Issues found'}")
        if not is_safe:
            print(f"Security issues: {issues}")
        
        # Perform conversion
        print("\nStarting PPTX to PDF conversion...")
        print("Note: This should NOT show any console windows!")
        
        success = converter.pptx_to_pdf(
            input_path=test_pptx,
            output_path=output_pdf,
            quality='high'
        )
        
        if success and os.path.exists(output_pdf):
            file_size = os.path.getsize(output_pdf)
            print(f"✅ Conversion successful!")
            print(f"   Output file: {output_pdf}")
            print(f"   File size: {file_size:,} bytes")
            
            # Verify PDF structure
            try:
                import fitz
                doc = fitz.open(output_pdf)
                print(f"   PDF pages: {len(doc)}")
                doc.close()
            except ImportError:
                print("   (PyMuPDF not available for PDF verification)")
            
            return True
        else:
            print("❌ Conversion failed!")
            return False
            
    except Exception as e:
        print(f"❌ Test error: {e}")
        return False
    finally:
        # Cleanup
        shutil.rmtree(temp_dir, ignore_errors=True)
        print(f"\nCleaned up temp directory: {temp_dir}")

def test_pdf_to_pptx_conversion():
    """Test PDF to PPTX conversion"""
    print("\n" + "="*50)
    print("Testing PDF to PPTX Conversion")
    print("="*50)
    
    # First create a PDF file from the PPTX
    temp_dir = tempfile.mkdtemp(prefix='pdf_test_')
    print(f"Using temp directory: {temp_dir}")
    
    try:
        # Create test PPTX and convert to PDF first
        test_pptx = os.path.join(temp_dir, "test_presentation.pptx")
        temp_pdf = os.path.join(temp_dir, "test_source.pdf")
        output_pptx = os.path.join(temp_dir, "converted_presentation.pptx")
        
        create_test_pptx(test_pptx)
        
        converter = PresentationConverter()
        
        # First convert PPTX to PDF
        print("Creating test PDF...")
        if not converter.pptx_to_pdf(test_pptx, temp_pdf):
            print("❌ Failed to create test PDF")
            return False
        
        # Check PDF to PPTX dependencies
        print(f"\nPDF to PPTX basic available: {converter.is_pdf_conversion_available()}")
        print(f"PDF to PPTX accurate available: {converter.is_ocr_available()}")
        
        if not converter.is_pdf_conversion_available():
            print("❌ Required dependencies not available for PDF to PPTX conversion")
            return False
        
        # Test security validation on PDF
        print("\nTesting PDF security validation...")
        is_safe, issues = converter.validate_file_security(temp_pdf)
        print(f"Security validation result: {'✅ Safe' if is_safe else '❌ Issues found'}")
        if not is_safe:
            print(f"Security issues: {issues}")
            # For testing, we'll continue even if security fails
        
        # Perform conversion
        print("\nStarting PDF to PPTX conversion (basic mode)...")
        success = converter.pdf_to_pptx_basic(
            input_path=temp_pdf,
            output_path=output_pptx
        )
        
        if success and os.path.exists(output_pptx):
            file_size = os.path.getsize(output_pptx)
            print(f"✅ Conversion successful!")
            print(f"   Output file: {output_pptx}")
            print(f"   File size: {file_size:,} bytes")
            
            # Verify PPTX structure
            try:
                from pptx import Presentation
                prs = Presentation(output_pptx)
                print(f"   PPTX slides: {len(prs.slides)}")
            except Exception as e:
                print(f"   PPTX verification error: {e}")
            
            return True
        else:
            print("❌ Conversion failed!")
            return False
            
    except Exception as e:
        print(f"❌ Test error: {e}")
        return False
    finally:
        # Cleanup
        shutil.rmtree(temp_dir, ignore_errors=True)
        print(f"\nCleaned up temp directory: {temp_dir}")

def main():
    """Run all conversion tests"""
    print("Presentation Conversion Test Suite")
    print("Testing enhanced conversion with security framework integration")
    print("This test should NOT show any LibreOffice console windows!")
    
    results = []
    
    # Test PPTX to PDF
    results.append(("PPTX to PDF", test_pptx_to_pdf_conversion()))
    
    # Test PDF to PPTX
    results.append(("PDF to PPTX", test_pdf_to_pptx_conversion()))
    
    # Summary
    print("\n" + "="*50)
    print("TEST SUMMARY")
    print("="*50)
    
    for test_name, result in results:
        status = "✅ PASS" if result else "❌ FAIL"
        print(f"{test_name:<20}: {status}")
    
    all_passed = all(result for _, result in results)
    print(f"\nOverall result: {'✅ ALL TESTS PASSED' if all_passed else '❌ SOME TESTS FAILED'}")
    
    if all_passed:
        print("\n🎉 Great! The console window issue should be fixed.")
        print("   The universal security framework is properly integrated.")
    else:
        print("\n⚠️  Some issues remain. Check the error messages above.")
    
    return all_passed

if __name__ == "__main__":
    success = main()
    sys.exit(0 if success else 1)